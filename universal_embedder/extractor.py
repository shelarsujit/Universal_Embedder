import os
import fitz  # PyMuPDF
from pytube import YouTube
from bs4 import BeautifulSoup
from urllib.request import urlopen
import pptx  # for pptx files
import docx  # for docx files
import zipfile
from .utils import generate_metadata
from .embedder import generate_embeddings
from .vectordb import store_embedding_in_vectordb
import traceback

def extract_text_from_pdf(pdf_path):
    try:
        text = ""
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        traceback.print_exc()
        return ""

def extract_text_from_pptx(pptx_path):
    try:
        text = ""
        presentation = pptx.Presentation(pptx_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        traceback.print_exc()
        return ""

def extract_text_from_docx(docx_path):
    try:
        text = ""
        doc = docx.Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + " "
        return text
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        traceback.print_exc()
        return ""

def extract_text_from_html(html_path):
    try:
        with open(html_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
        return soup.get_text()
    except Exception as e:
        print(f"Error extracting text from HTML: {e}")
        traceback.print_exc()
        return ""

def extract_text_from_idml(idml_path):
    try:
        text = ""
        with zipfile.ZipFile(idml_path, 'r') as zip_ref:
            for filename in zip_ref.namelist():
                if filename.endswith(".xml"):
                    with zip_ref.open(filename) as file:
                        soup = BeautifulSoup(file, 'xml')
                        text += soup.get_text() + " "
        return text
    except Exception as e:
        print(f"Error extracting text from IDML: {e}")
        traceback.print_exc()
        return ""

def extract_text_from_url(url):
    try:
        page = urlopen(url)
        soup = BeautifulSoup(page, 'html.parser')
        return soup.get_text()
    except Exception as e:
        print(f"Error extracting text from URL: {e}")
        traceback.print_exc()
        return ""

def extract_audio_from_youtube(link):
    try:
        yt = YouTube(link)
        # Download only audio
        audio_stream = yt.streams.filter(only_audio=True).first()
        download_path = audio_stream.download()
        return download_path
    except Exception as e:
        print(f"Error extracting audio from YouTube: {e}")
        traceback.print_exc()
        return ""

def process_file(file_path):
    metadata = generate_metadata(file_path)
    file_extension = os.path.splitext(file_path)[-1].lower()
    text = ""
    if file_extension == '.pdf':
        text = extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        text = extract_text_from_pptx(file_path)
    elif file_extension == '.docx':
        text = extract_text_from_docx(file_path)
    elif file_extension == '.html':
        text = extract_text_from_html(file_path)
    elif file_extension == '.idml':
        text = extract_text_from_idml(file_path)
    else:
        print(f"Unsupported file format: {file_extension}")
        return

    if text:
        embedding = generate_embeddings(text)
        store_embedding_in_vectordb(embedding, metadata)

def process_link(link):
    metadata = generate_metadata(link)
    if 'youtube.com' in link or 'youtu.be' in link:
        audio_path = extract_audio_from_youtube(link)
        # Further audio transcription can be done here
        print(f"Downloaded audio to: {audio_path}")
    else:
        text = extract_text_from_url(link)
        if text:
            embedding = generate_embeddings(text)
            store_embedding_in_vectordb(embedding, metadata)
